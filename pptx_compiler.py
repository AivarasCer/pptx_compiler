import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


def images_to_ppt(folder_path, output_file="output_presentation.pptx"):
    if not os.path.exists(folder_path):
        print("The specified folder does not exist.")
        return

    presentation = Presentation()

    supported_formats = ('.jpg', '.jpeg', '.png', '.gif')

    for filename in os.listdir(folder_path):
        if filename.lower().endswith(supported_formats):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            image_path = os.path.join(folder_path, filename)

            try:
                with Image.open(image_path) as img:
                    width, height = img.size
                    dpi = img.info.get('dpi', (96, 96))
                    width_inches = width / dpi[0]
                    height_inches = height / dpi[1]

                slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(width_inches), height=Inches(height_inches))
            except Exception as e:
                print(f"Could not add {filename}: {e}")

    presentation.save(output_file)
    print(f"PowerPoint presentation saved as {output_file}")

folder_path = input("Enter the folder path containing the images: ")
output_file = input("Enter the output PowerPoint file name (default: output_presentation.pptx): ") or "output_presentation.pptx"

images_to_ppt(folder_path, output_file)
